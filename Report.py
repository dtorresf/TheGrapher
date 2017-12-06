import Methods
import Config
import Pptxr
from pptx import Presentation
from datetime import datetime
import pandas as pd 

class Report:
	'''Class that represents a report'''
	def __init__(self):
		self.date = ''
		self.name = ''
		self.title = ''
		self.servers = list()
		self.switches = list()
		self.zfs = list()
		self.cf = Config.Config()

	def loadreport(self,date,name,title,servers,switches,zfs,cf):
		self.name = date
		self.name = name
		self.title = title
		self.servers = servers
		self.switches = switches
		self.zfs = zfs
		self.cf = cf
	def loadreport2(self,date,name,title,servers,switches,cf):
		self.name = date
		self.name = name
		self.title = title
		self.servers = servers
		self.switches = switches
		self.zfs = ''
		self.cf = cf

	def generatepptxreport(self):
		date = datetime.now().strftime("%d%m%Y-%H%M%S")
		report_name = self.cf.variables['pptx_report'] + '/' + 'Exalogic_Report_'+ date + '.pptx'
		#Title presentation
		pptxr = Pptxr.Pptxr(self.cf.variables['pptx_template'])
		title_slide = pptxr.prs.slides[0]
		title = title_slide.shapes.title
		title.text = self.title
		placeholder_content = title_slide.placeholders[1]
		placeholder_content.text = 'REPORT - ' + date 
		#Servers
		for s in self.servers:
			s.serverpptxreport(self.cf,pptxr)
		#Switches
		for sw in self.switches:
			sw.switchpptxreport(self.cf,pptxr)
		#ZFS
		self.zfs.zfspptxreport(self.cf,pptxr)

		pptxr.prs.save(self.name)

	def generatepptxreport2(self):
		date = datetime.now().strftime("%d%m%Y-%H%M%S")
		report_name = self.cf.variables['pptx_report'] + '/' + 'Exalogic_Report_'+ date + '.pptx'
		#Title presentation
		pptxr = Pptxr.Pptxr(self.cf.variables['pptx_template'])
		title_slide = pptxr.prs.slides[0]
		title = title_slide.shapes.title
		title.text = self.title
		placeholder_content = title_slide.placeholders[1]
		placeholder_content.text = 'REPORT - ' + date 
		#Servers
		for s in self.servers:
			s.serverpptxreport(self.cf,pptxr)
		#Switches
		for sw in self.switches:
			sw.switchpptxreport(self.cf,pptxr)
		#ZFS
		#self.zfs.zfspptxreport(self.cf,pptxr)

		pptxr.prs.save(self.name)
