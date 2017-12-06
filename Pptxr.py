import Methods
import Config
import pandas as pd 
from pptx import Presentation
from datetime import datetime

class Pptxr:
	'''Class that represents a Server with atributtes to graph'''
	def __init__(self,template):
		self.prs = Presentation(template)
		self.template = template

	def twoimageslidepptx(self,titletext,img1,img2,textimg1,textimg2):
		image_slide = self.prs.slides.add_slide(self.prs.slide_layouts[25])
		title = image_slide.shapes.title
		title.text = titletext
			
		placeholder = image_slide.placeholders[1] #Capture first image placeholder for CPU
		picture = placeholder.insert_picture(img1)
		image_slide.placeholders[2].text = textimg1
			
		placeholder = image_slide.placeholders[13]  # idx key, not position
		picture = placeholder.insert_picture(img2)
		image_slide.placeholders[14].text = textimg2

	def oneimageslidepptx(self,titletext,img,textimg):
		image_slide = self.prs.slides.add_slide(self.prs.slide_layouts[24])
		title = image_slide.shapes.title
		title.text = titletext

		placeholder = image_slide.placeholders[1] #Capture first image placeholder for TX
		picture = placeholder.insert_picture(img)
		image_slide.placeholders[2].text = textimg